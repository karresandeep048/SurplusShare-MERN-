import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6] # Blank layout

    # Color Palette
    DARK_BG = RGBColor(15, 23, 42)       # Slate 900
    LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50
    CARD_BG = RGBColor(255, 255, 255)    # Pure White
    CARD_BORDER = RGBColor(226, 232, 240) # Slate 200
    
    PRIMARY_GREEN = RGBColor(16, 185, 129)  # Emerald 500
    DARK_GREEN = RGBColor(5, 150, 105)      # Emerald 600
    LIGHT_GREEN = RGBColor(209, 250, 229)   # Emerald 100
    TEAL_ACCENT = RGBColor(13, 148, 136)    # Teal 600
    
    TEXT_DARK = RGBColor(15, 23, 42)        # Slate 900
    TEXT_MUTED = RGBColor(100, 116, 139)    # Slate 500
    TEXT_LIGHT = RGBColor(241, 245, 249)    # Slate 100
    ACCENT_ROSE = RGBColor(244, 63, 94)     # Rose 500
    ACCENT_BLUE = RGBColor(59, 130, 246)    # Blue 500
    ACCENT_AMBER = RGBColor(217, 119, 6)    # Amber 600

    def add_header(slide, title_text, category_text=None, is_dark=False):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.1))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        if category_text:
            p_cat = tf.paragraphs[0]
            p_cat.text = category_text.upper()
            p_cat.font.size = Pt(11)
            p_cat.font.bold = True
            p_cat.font.color.rgb = PRIMARY_GREEN if is_dark else DARK_GREEN
            
            p_title = tf.add_paragraph()
        else:
            p_title = tf.paragraphs[0]

        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_LIGHT if is_dark else TEXT_DARK

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BG
    bg1.line.fill.background()

    tbox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.0))
    tf1 = tbox.text_frame
    tf1.word_wrap = True
    
    p0 = tf1.paragraphs[0]
    p0.text = "🌱 MERN STACK SUSTAINABILITY PLATFORM"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY_GREEN
    
    p1 = tf1.add_paragraph()
    p1.text = "SurplusShare"
    p1.font.size = Pt(50)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    p1.space_before = Pt(12)

    p2 = tf1.add_paragraph()
    p2.text = "Save Food. Share Good. Connecting Surplus Meals with Communities in Real-Time."
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    p2.space_before = Pt(14)

    features_summary = [
        ("📍 Geospatial Matching", "Interactive Leaflet live food map"),
        ("📧 Live Email Center", "Instant digital claim passes & arrival alerts"),
        ("🔐 6-Digit Handover OTP", "Cryptographic pickup verification"),
        ("🌍 ESG Impact Metrics", "Calculated CO2, water & meal savings")
    ]
    
    for i, (f_title, f_desc) in enumerate(features_summary):
        card = add_card(slide1, Inches(1.0 + i * 2.88), Inches(5.4), Inches(2.7), Inches(1.3), 
                        bg_color=RGBColor(30, 41, 59), border_color=RGBColor(51, 65, 85))
        cbox = slide1.shapes.add_textbox(Inches(1.1 + i * 2.88), Inches(5.5), Inches(2.5), Inches(1.1))
        ctf = cbox.text_frame
        ctf.word_wrap = True
        cp1 = ctf.paragraphs[0]
        cp1.text = f_title
        cp1.font.size = Pt(13)
        cp1.font.bold = True
        cp1.font.color.rgb = PRIMARY_GREEN
        cp2 = ctf.add_paragraph()
        cp2.text = f_desc
        cp2.font.size = Pt(10)
        cp2.font.color.rgb = RGBColor(203, 213, 225)
        cp2.space_before = Pt(4)

    # ==========================================
    # SLIDE 2: Problem Statement & Market Gaps
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide2, "The Food Waste Paradox & Market Gaps", "01 / PROBLEM STATEMENT")

    problems = [
        ("Huge Quantities of Edible Waste", 
         "Restaurants, catering businesses, university canteens, and bakeries prepare fresh food daily, discarding high-quality surplus meals simply due to lack of immediate redistribution channels.", 
         "🚨 Widespread Food Insecurity", ACCENT_ROSE),
        ("Absence of Real-Time Infrastructure", 
         "Traditional food charity operates on batch logistics with long lead times, making same-day perishable food rescue nearly impossible before food expires.", 
         "⏱️ Zero Perishable Support", ACCENT_AMBER),
        ("Security & Communication Deficit", 
         "Uncoordinated pickups cause food safety concerns, unauthorized collections, lack of instant arrival alerts, and no proof of handover between donor and receiver.", 
         "🔒 Verification & Notification Gap", ACCENT_BLUE)
    ]

    for i, (title, desc, tag, color) in enumerate(problems):
        add_card(slide2, Inches(0.8 + i * 3.95), Inches(1.8), Inches(3.75), Inches(4.8))
        
        badge = add_card(slide2, Inches(1.1 + i * 3.95), Inches(2.1), Inches(3.15), Inches(0.55), 
                         bg_color=RGBColor(241, 245, 249), border_color=color)
        b_box = slide2.shapes.add_textbox(Inches(1.1 + i * 3.95), Inches(2.15), Inches(3.15), Inches(0.45))
        b_tf = b_box.text_frame
        b_p = b_tf.paragraphs[0]
        b_p.text = tag
        b_p.font.size = Pt(12)
        b_p.font.bold = True
        b_p.font.color.rgb = color
        b_p.alignment = PP_ALIGN.CENTER
        
        box = slide2.shapes.add_textbox(Inches(1.0 + i * 3.95), Inches(2.9), Inches(3.35), Inches(3.5))
        tf = box.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(17)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(12)

    # ==========================================
    # SLIDE 3: Proposed Solution
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide3, "SurplusShare: The Direct Hyperlocal Food Bridge", "02 / PROPOSED SOLUTION")

    solution_pillars = [
        ("Donor Onboarding & Instant Postings", "Suppliers easily post surplus items in seconds with quantity, expiry timer, dietary labels, pickup windows, and geo-coordinates."),
        ("Map-Based Community Discovery", "Receivers browse nearby fresh & urgent food via interactive Leaflet map, radius filters, and dietary preference filters."),
        ("Atomic Claims & Quantity Management", "Instant reservation deduction prevents overselling; partial claim allows multiple community members to share bulk batches."),
        ("Interactive Email Notification Center", "Instant email claim passes with 6-digit OTP, resend option on UI, and instant arrival alerts to donors."),
        ("Two-Way Pickup Verification & ESG Impact", "Secure 6-digit OTP ensures authenticated handover while computing CO2, water, and meal rescue analytics in real-time.")
    ]

    for i, (title, desc) in enumerate(solution_pillars):
        add_card(slide3, Inches(0.8), Inches(1.8 + i * 1.02), Inches(11.7), Inches(0.88), 
                 bg_color=CARD_BG, border_color=CARD_BORDER)
        
        add_card(slide3, Inches(1.0), Inches(1.95 + i * 1.02), Inches(0.6), Inches(0.55), 
                 bg_color=LIGHT_GREEN, border_color=PRIMARY_GREEN)
        num_box = slide3.shapes.add_textbox(Inches(1.0), Inches(1.98 + i * 1.02), Inches(0.6), Inches(0.5))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = f"0{i+1}"
        np.font.size = Pt(13)
        np.font.bold = True
        np.font.color.rgb = DARK_GREEN
        np.alignment = PP_ALIGN.CENTER
        
        tbox = slide3.shapes.add_textbox(Inches(1.8), Inches(1.88 + i * 1.02), Inches(10.5), Inches(0.7))
        ttf = tbox.text_frame
        ttf.word_wrap = True
        tp1 = ttf.paragraphs[0]
        tp1.text = title
        tp1.font.size = Pt(14)
        tp1.font.bold = True
        tp1.font.color.rgb = TEXT_DARK
        
        tp2 = ttf.add_paragraph()
        tp2.text = desc
        tp2.font.size = Pt(11)
        tp2.font.color.rgb = TEXT_MUTED
        tp2.space_before = Pt(2)

    # ==========================================
    # SLIDE 4: User Roles & Workflows
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide4, "Target Stakeholders & Dual-Sided Workflows", "03 / USER ROLES")

    roles = [
        ("🍽️ Food Suppliers / Donors", 
         ["Restaurants, Bakeries, Canteens, Caterers",
          "Create food listings with photo, portions & time",
          "Receive automated reservation & arrival emails",
          "Enter 6-digit receiver OTP to confirm handover",
          "View donor impact analytics & meals rescued"],
         DARK_GREEN, LIGHT_GREEN),
        ("🤝 Food Receivers / Community", 
         ["Individuals, NGOs, Volunteers, Shelters",
          "Search food by map location, urgency & diet",
          "Reserve meals & trigger instant email passes",
          "Track pickup route & send 1-click arrival alert",
          "Present unique 6-digit OTP during collection"],
         RGBColor(2, 132, 199), RGBColor(224, 242, 254)),
        ("🛡️ System Administrators", 
         ["Platform Governance & Community Trust",
          "Supervise listings, reservations & users",
          "Automated background auto-expiry management",
          "Ecosystem sustainability & city-wide metrics",
          "Test & monitor live email dispatch pipeline"],
         RGBColor(124, 58, 237), RGBColor(243, 232, 255))
    ]

    for i, (title, items, primary_col, bg_light) in enumerate(roles):
        add_card(slide4, Inches(0.8 + i * 3.95), Inches(1.8), Inches(3.75), Inches(4.8))
        
        add_card(slide4, Inches(0.8 + i * 3.95), Inches(1.8), Inches(3.75), Inches(0.9), 
                 bg_color=bg_light, border_color=primary_col)
        h_box = slide4.shapes.add_textbox(Inches(0.9 + i * 3.95), Inches(1.95), Inches(3.55), Inches(0.6))
        htf = h_box.text_frame
        htf.word_wrap = True
        hp = htf.paragraphs[0]
        hp.text = title
        hp.font.size = Pt(14)
        hp.font.bold = True
        hp.font.color.rgb = primary_col

        l_box = slide4.shapes.add_textbox(Inches(1.0 + i * 3.95), Inches(2.85), Inches(3.35), Inches(3.6))
        ltf = l_box.text_frame
        ltf.word_wrap = True
        for j, item in enumerate(items):
            p = ltf.add_paragraph() if j > 0 else ltf.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.space_before = Pt(8)

    # ==========================================
    # SLIDE 5: Technical Architecture (MERN Stack)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide5, "Full-Stack System Architecture", "04 / TECHNICAL ARCHITECTURE")

    tech_boxes = [
        ("Client Tier (React 18 + Vite)", 
         ["• Tailwind CSS v4 & Lucide Icons for responsive UX",
          "• React-Leaflet for geospatial map visualization",
          "• Email Sending & Delivery Status controls in UI",
          "• Real-time Order Tracker with arrival polling",
          "• Dynamic ESG Impact Calculator component"],
         Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.8)),

        ("Application Tier (Express 5 + Node)", 
         ["• RESTful API Architecture with MVC separation",
          "• Auth Middleware with JWT signature verification",
          "• Upgraded Nodemailer Gmail notification pipeline",
          "• Automated autoExpireItems() lifecycle engine",
          "• Dedicated /api/reservations/resend-email endpoints"],
         Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.8)),

        ("Data Tier (MongoDB Atlas)", 
         ["• User Schema: Roles, preferences & impact",
          "• FoodListing Schema: Geolocation & availability",
          "• Reservation Schema: Timestamps & pickup codes",
          "• Idempotent database seeding scripts",
          "• Atomic transactions & quantity decrement"],
         Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8))
    ]

    for title, points, left, top, w, h in tech_boxes:
        add_card(slide5, left, top, w, h)
        
        add_card(slide5, left, top, w, Inches(0.8), bg_color=RGBColor(241, 245, 249), border_color=CARD_BORDER)
        t_box = slide5.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), w - Inches(0.3), Inches(0.5))
        ttf = t_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(13)
        tp.font.bold = True
        tp.font.color.rgb = DARK_GREEN
        
        p_box = slide5.shapes.add_textbox(left + Inches(0.2), top + Inches(0.9), w - Inches(0.4), h - Inches(1.0))
        ptf = p_box.text_frame
        ptf.word_wrap = True
        for j, pt_str in enumerate(points):
            p = ptf.add_paragraph() if j > 0 else ptf.paragraphs[0]
            p.text = pt_str
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.space_before = Pt(10)

    # ==========================================
    # SLIDE 6: Dedicated Email Upgradation & Notification System
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide6, "Upgraded Email & Automated Notification Pipeline", "05 / EMAIL NOTIFICATION SYSTEM")

    email_features = [
        ("1. Instant Receiver Claim Pass", 
         "When a user reserves food, an automated HTML email pass is dispatched containing the food name, donor address, pickup window, and unique 6-digit OTP code.",
         DARK_GREEN),
        ("2. Real-Time Donor Claim Alert", 
         "Suppliers instantly receive an email alert with recipient details, portion count, and safety instructions when their surplus food is claimed.",
         RGBColor(2, 132, 199)),
        ("3. On-Demand 'Email Pass' Option in UI", 
         "Users can click 'Email Pass' on MyReservations, Order Tracker, and Dashboard to send/resend their vouchers to any email address anytime.",
         ACCENT_AMBER),
        ("4. Resilient Fallback & Diagnostics", 
         "Configured with Nodemailer Gmail SMTP (16-char App Passwords), 5-second socket timeouts, and self-healing formatted console logging.",
         RGBColor(124, 58, 237))
    ]

    for i, (title, desc, col) in enumerate(email_features):
        col_idx = i % 2
        row_idx = i // 2
        l = Inches(0.8 + col_idx * 6.0)
        t = Inches(1.8 + row_idx * 2.5)
        
        add_card(slide6, l, t, Inches(5.7), Inches(2.2))
        add_card(slide6, l, t, Inches(0.15), Inches(2.2), bg_color=col, border_color=col)
        
        box = slide6.shapes.add_textbox(l + Inches(0.35), t + Inches(0.2), Inches(5.15), Inches(1.8))
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_DARK
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    # ==========================================
    # SLIDE 7: Key Features & Innovation Highlights
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide7, "Key Innovation Highlights & Capabilities", "06 / CORE FEATURES")

    features = [
        ("📍 Interactive Geospatial Map", "Visualizes nearby food pins (Urgent ⚡ vs Fresh 🌱), allows radius filtering (5-50km), and enables one-click directions.", Inches(0.8), Inches(1.8)),
        ("🔐 6-Digit Pickup Code Verification", "Zero-friction verification system: Receivers provide a unique 6-digit OTP; suppliers verify it on the terminal to complete handover.", Inches(6.8), Inches(1.8)),
        ("🛵 Live Order & Arrival Tracker", "Live progress route simulation with one-touch 'Notify Donor I Have Arrived' sending instant email & in-app alerts.", Inches(0.8), Inches(3.5)),
        ("📧 Integrated Email Center & Passes", "Direct 'Email Pass' and 'Resend Pass' buttons in the UI dispatching full HTML vouchers to user inboxes.", Inches(6.8), Inches(3.5)),
        ("⏳ Self-Healing Auto-Expiry", "Automated cron/service marks past-due food as EXPIRED and releases abandoned holds back to the community pool.", Inches(0.8), Inches(5.2)),
        ("📊 ESG Impact Analytics", "Live dashboard computing total meals rescued, CO2 avoided (2.5kg/meal), and water preserved (140L/meal).", Inches(6.8), Inches(5.2))
    ]

    for title, desc, l, t in features:
        add_card(slide7, l, t, Inches(5.7), Inches(1.45))
        box = slide7.shapes.add_textbox(l + Inches(0.2), t + Inches(0.15), Inches(5.3), Inches(1.15))
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = DARK_GREEN
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(4)

    # ==========================================
    # SLIDE 8: Database Design & Data Models
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide8, "Database Schema & Entity Relationships", "07 / DATA DESIGN")

    schemas = [
        ("User Model", 
         ["name: String (Required)",
          "email: String (Unique, Indexed)",
          "password: String (Bcrypt Hashed)",
          "role: ['supplier', 'receiver', 'admin']",
          "dietaryPreferences: [String]",
          "mealsRescued: Number (Default: 0)",
          "location: String, profileImage: String"]),
        
        ("FoodListing Model", 
         ["supplier: ObjectId (Ref: 'User')",
          "foodName, description, unit: String",
          "quantity, availableQuantity: Number",
          "foodType: ['Vegetarian', 'Vegan', ...]",
          "pickupStart, pickupEnd, expiryTime: Date",
          "coordinates: { lat: Number, lng: Number }",
          "status: ['AVAILABLE', 'RESERVED', 'COLLECTED', 'EXPIRED']"]),

        ("Reservation Model", 
         ["foodListing: ObjectId (Ref: 'FoodListing')",
          "receiver: ObjectId (Ref: 'User')",
          "quantity: Number (Claimed count)",
          "pickupCode: String (6-digit OTP)",
          "pickerArrived: Boolean (Arrival ping)",
          "status: ['RESERVED', 'COLLECTED', 'CANCELLED', 'EXPIRED']",
          "reservedAt, arrivedAt, collectedAt: Date"])
    ]

    for i, (name, fields) in enumerate(schemas):
        add_card(slide8, Inches(0.8 + i * 3.95), Inches(1.8), Inches(3.75), Inches(4.8))
        
        add_card(slide8, Inches(0.8 + i * 3.95), Inches(1.8), Inches(3.75), Inches(0.65), 
                 bg_color=DARK_BG, border_color=DARK_BG)
        h_box = slide8.shapes.add_textbox(Inches(0.9 + i * 3.95), Inches(1.9), Inches(3.55), Inches(0.45))
        htf = h_box.text_frame
        hp = htf.paragraphs[0]
        hp.text = f"📦 {name}"
        hp.font.size = Pt(13)
        hp.font.bold = True
        hp.font.color.rgb = PRIMARY_GREEN

        f_box = slide8.shapes.add_textbox(Inches(0.95 + i * 3.95), Inches(2.6), Inches(3.45), Inches(3.8))
        ftf = f_box.text_frame
        ftf.word_wrap = True
        for j, f in enumerate(fields):
            p = ftf.add_paragraph() if j > 0 else ftf.paragraphs[0]
            p.text = f"• {f}"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
            p.space_before = Pt(6)

    # ==========================================
    # SLIDE 9: End-to-End User Experience & Flow
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide9, "End-to-End Operational Lifecycle", "08 / USER JOURNEY")

    steps = [
        ("1. Food Listed", "Supplier inputs food details, quantities, photo & expiry window.", DARK_GREEN),
        ("2. Real-Time Discovery", "Receivers explore live pins & filter by distance or diet.", RGBColor(2, 132, 199)),
        ("3. Claim & Email Pass", "6-digit OTP generated & emailed with on-demand resend option.", RGBColor(217, 119, 6)),
        ("4. Arrival Ping", "Receiver reaches venue & clicks 'Notify Arrival'; donor alerted by email.", RGBColor(124, 58, 237)),
        ("5. Code Handover", "Donor validates 6-digit OTP; meal marked COLLECTED & logged.", DARK_GREEN)
    ]

    for i, (stitle, sdesc, col) in enumerate(steps):
        add_card(slide9, Inches(0.8 + i * 2.38), Inches(2.0), Inches(2.2), Inches(4.5))
        
        add_card(slide9, Inches(1.4 + i * 2.38), Inches(2.3), Inches(1.0), Inches(0.8), 
                 bg_color=RGBColor(241, 245, 249), border_color=col)
        c_box = slide9.shapes.add_textbox(Inches(1.4 + i * 2.38), Inches(2.45), Inches(1.0), Inches(0.5))
        ctf = c_box.text_frame
        cp = ctf.paragraphs[0]
        cp.text = f"0{i+1}"
        cp.font.size = Pt(18)
        cp.font.bold = True
        cp.font.color.rgb = col
        cp.alignment = PP_ALIGN.CENTER

        t_box = slide9.shapes.add_textbox(Inches(0.9 + i * 2.38), Inches(3.4), Inches(2.0), Inches(2.8))
        ttf = t_box.text_frame
        ttf.word_wrap = True
        tp1 = ttf.paragraphs[0]
        tp1.text = stitle
        tp1.font.size = Pt(13)
        tp1.font.bold = True
        tp1.font.color.rgb = TEXT_DARK
        tp1.alignment = PP_ALIGN.CENTER

        tp2 = ttf.add_paragraph()
        tp2.text = sdesc
        tp2.font.size = Pt(11)
        tp2.font.color.rgb = TEXT_MUTED
        tp2.space_before = Pt(8)
        tp2.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 10: Environmental & Social Impact (ESG)
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide10, "Measurable Environmental & Social Impact", "09 / IMPACT & SUSTAINABILITY")

    impact_stats = [
        ("🍲 1,420+", "Meals Rescued This Month", "Redirected to local community members & shelters across Bengaluru.", DARK_GREEN),
        ("🌿 3,550 kg", "CO2 Emissions Prevented", "Direct landfill greenhouse gas reduction (estimated at 2.5kg CO2e per meal).", RGBColor(13, 148, 136)),
        ("💧 198,800 L", "Water Footprint Preserved", "Virtual water footprint saved through upstream food salvage (140L/meal).", RGBColor(2, 132, 199)),
        ("💰 ₹1,70,000+", "Community Value Generated", "Direct economic relief delivered to underserved families and partner NGOs.", RGBColor(217, 119, 6))
    ]

    for i, (val, title, desc, col) in enumerate(impact_stats):
        col_idx = i % 2
        row_idx = i // 2
        l = Inches(0.8 + col_idx * 6.0)
        t = Inches(1.8 + row_idx * 2.5)
        
        add_card(slide10, l, t, Inches(5.7), Inches(2.2), bg_color=RGBColor(248, 250, 252))
        
        box = slide10.shapes.add_textbox(l + Inches(0.3), t + Inches(0.2), Inches(5.1), Inches(1.8))
        tf = box.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = val
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_DARK
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = TEXT_MUTED
        p3.space_before = Pt(4)

    # ==========================================
    # SLIDE 11: Future Roadmap & Scalability
    # ==========================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide11, "Strategic Roadmap & Future Evolution", "10 / FUTURE SCOPE")

    roadmap_items = [
        ("Phase 1: Logistics & Volunteer Network", 
         "Introduce crowdsourced volunteer pickup drivers for bulk food donations where receivers lack immediate transportation.", 
         "Next Quarter"),
        ("Phase 2: AI Expiry & Surplus Prediction", 
         "ML-based forecasting model advising restaurants on daily preparation volumes and auto-boosting urgent food before expiration.", 
         "Mid Term"),
        ("Phase 3: QR Code & Mobile PWA", 
         "Native mobile PWA with camera QR code scanning for instantaneous supplier-receiver handovers and offline sync.", 
         "Mid Term"),
        ("Phase 4: Corporate CSR & Tax Incentive Portal", 
         "Automated tax benefit & CSR impact certificates for corporate canteens and hotel chains donating excess food.", 
         "Long Term")
    ]

    for i, (title, desc, time_tag) in enumerate(roadmap_items):
        add_card(slide11, Inches(0.8), Inches(1.8 + i * 1.25), Inches(11.7), Inches(1.1))
        
        add_card(slide11, Inches(1.0), Inches(2.0 + i * 1.25), Inches(1.6), Inches(0.55), 
                 bg_color=LIGHT_GREEN, border_color=PRIMARY_GREEN)
        tb_box = slide11.shapes.add_textbox(Inches(1.0), Inches(2.08 + i * 1.25), Inches(1.6), Inches(0.45))
        tbtf = tb_box.text_frame
        tbp = tbtf.paragraphs[0]
        tbp.text = time_tag
        tbp.font.size = Pt(11)
        tbp.font.bold = True
        tbp.font.color.rgb = DARK_GREEN
        tbp.alignment = PP_ALIGN.CENTER

        c_box = slide11.shapes.add_textbox(Inches(2.8), Inches(1.9 + i * 1.25), Inches(9.5), Inches(0.9))
        ctf = c_box.text_frame
        ctf.word_wrap = True
        cp1 = ctf.paragraphs[0]
        cp1.text = title
        cp1.font.size = Pt(14)
        cp1.font.bold = True
        cp1.font.color.rgb = TEXT_DARK
        
        cp2 = ctf.add_paragraph()
        cp2.text = desc
        cp2.font.size = Pt(11)
        cp2.font.color.rgb = TEXT_MUTED
        cp2.space_before = Pt(3)

    # ==========================================
    # SLIDE 12: Conclusion & Q&A (Dark Theme)
    # ==========================================
    slide12 = prs.slides.add_slide(blank_slide_layout)
    bg12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = DARK_BG
    bg12.line.fill.background()

    tbox12 = slide12.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
    tf12 = tbox12.text_frame
    tf12.word_wrap = True

    p0 = tf12.paragraphs[0]
    p0.text = "🌱 TRANSFORMATIVE COMMUNITY SUSTAINABILITY"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY_GREEN

    p1 = tf12.add_paragraph()
    p1.text = "Thank You!"
    p1.font.size = Pt(50)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    p1.space_before = Pt(8)

    p2 = tf12.add_paragraph()
    p2.text = "SurplusShare turns food waste into community nourishment through thoughtful technology."
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    p2.space_before = Pt(12)

    add_card(slide12, Inches(1.0), Inches(4.5), Inches(11.3), Inches(2.0), 
             bg_color=RGBColor(30, 41, 59), border_color=RGBColor(51, 65, 85))
    
    cred_box = slide12.shapes.add_textbox(Inches(1.2), Inches(4.65), Inches(10.9), Inches(1.7))
    ctf12 = cred_box.text_frame
    ctf12.word_wrap = True
    
    cp_title = ctf12.paragraphs[0]
    cp_title.text = "🚀 LIVE DEMO CREDENTIALS (Default Password: password123)"
    cp_title.font.size = Pt(13)
    cp_title.font.bold = True
    cp_title.font.color.rgb = PRIMARY_GREEN

    cp_s = ctf12.add_paragraph()
    cp_s.text = "• Supplier: demo.supplier@surplusshare.com  |  Green Bowl Restaurant (Post & Verify)"
    cp_s.font.size = Pt(12)
    cp_s.font.color.rgb = RGBColor(226, 232, 240)
    cp_s.space_before = Pt(6)

    cp_r = ctf12.add_paragraph()
    cp_r.text = "• Receiver: demo.receiver@surplusshare.com  |  Arjun / Community Member (Find, Reserve & Track)"
    cp_r.font.size = Pt(12)
    cp_r.font.color.rgb = RGBColor(226, 232, 240)
    cp_r.space_before = Pt(4)

    cp_a = ctf12.add_paragraph()
    cp_a.text = "• Admin: demo.admin@surplusshare.com  |  System & Ecosystem Management"
    cp_a.font.size = Pt(12)
    cp_a.font.color.rgb = RGBColor(226, 232, 240)
    cp_a.space_before = Pt(4)

    output_path = "/Users/apple/Desktop/SurplusShare/SurplusShare_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation successfully updated and saved to {output_path}")

if __name__ == '__main__':
    create_presentation()
